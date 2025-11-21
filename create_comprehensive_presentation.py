#!/usr/bin/env python3
"""
Comprehensive API Automation Framework Presentation Generator
Creates a PowerPoint presentation covering:
- Problem Statement
- Solution Overview
- Architecture Structure
- Key Features
- Demo
- Benefits
- Challenges and Learning
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    # Style the title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_section_header(prs, title):
    """Add a section header slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section header layout
    title_shape = slide.shapes.title
    title_shape.text = title

    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]

    title_shape.text = title

    # Style title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add content
    text_frame = content_shape.text_frame
    text_frame.clear()

    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = item['text']
        p.level = item.get('level', 0)
        p.font.size = Pt(item.get('size', 18))

        if item.get('bold', False):
            p.font.bold = True

    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.25), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    for i, item in enumerate(left_items):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item['text']
        p.level = item.get('level', 0)
        p.font.size = Pt(item.get('size', 16))
        if item.get('bold', False):
            p.font.bold = True

    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.25), Inches(1.5), Inches(4.25), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    for i, item in enumerate(right_items):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item['text']
        p.level = item.get('level', 0)
        p.font.size = Pt(item.get('size', 16))
        if item.get('bold', False):
            p.font.bold = True

    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "API Automation Framework",
        "Comprehensive Testing Solution for Boundary Management & Localization Services"
    )

    # Slide 2: Problem Statement
    add_section_header(prs, "Problem Statement")

    add_content_slide(prs, "The Challenge", [
        {'text': 'Manual API testing is time-consuming and error-prone', 'size': 20, 'bold': True},
        {'text': 'Complex multi-step workflows requiring sequential execution', 'level': 1, 'size': 18},
        {'text': 'Need to validate 15+ interconnected API endpoints', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Boundary management complexities', 'size': 20, 'bold': True},
        {'text': '7-level hierarchical boundary structures', 'level': 1, 'size': 18},
        {'text': 'Multi-language localization support (English, French, Portuguese)', 'level': 1, 'size': 18},
        {'text': 'Template generation, file operations, and data processing', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Lack of comprehensive test reporting and traceability', 'size': 20, 'bold': True},
    ])

    # Slide 3: Solution Overview
    add_section_header(prs, "Solution Overview")

    add_content_slide(prs, "API Automation Framework", [
        {'text': 'Python-based automation framework using pytest', 'size': 20, 'bold': True},
        {'text': '', 'size': 18},
        {'text': 'What it does:', 'size': 20, 'bold': True},
        {'text': 'Automates end-to-end testing of boundary management APIs', 'level': 1, 'size': 18},
        {'text': 'Validates complete workflows from hierarchy creation to data processing', 'level': 1, 'size': 18},
        {'text': 'Supports template generation, file operations, and multi-language testing', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Key capabilities:', 'size': 20, 'bold': True},
        {'text': '15 sequential tests covering entire workflow', 'level': 1, 'size': 18},
        {'text': 'Automated template download, population, and upload', 'level': 1, 'size': 18},
        {'text': 'Dynamic hierarchy generation and validation', 'level': 1, 'size': 18},
        {'text': 'Comprehensive HTML and Allure reporting', 'level': 1, 'size': 18},
    ])

    # Slide 4: Architecture - High Level
    add_section_header(prs, "Architecture Structure")

    add_content_slide(prs, "System Architecture", [
        {'text': 'Modular Design Pattern', 'size': 22, 'bold': True},
        {'text': '', 'size': 18},
        {'text': 'Test Layer (tests/)', 'size': 20, 'bold': True},
        {'text': '15 sequential test modules', 'level': 1, 'size': 18},
        {'text': 'Boundary hierarchy, localization, file operations, data processing', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Utilities Layer (utils/)', 'size': 20, 'bold': True},
        {'text': 'api_client.py - HTTP client wrapper with auth', 'level': 1, 'size': 18},
        {'text': 'auth.py - OAuth2 token management', 'level': 1, 'size': 18},
        {'text': 'config.py - Environment configuration', 'level': 1, 'size': 18},
        {'text': 'data_loader.py - JSON payload loader', 'level': 1, 'size': 18},
        {'text': 'request_info.py - Request metadata builder', 'level': 1, 'size': 18},
    ])

    # Slide 5: Architecture - Components
    add_two_column_slide(
        prs,
        "Framework Components",
        [
            {'text': 'Data Layer', 'size': 20, 'bold': True},
            {'text': 'Payloads: JSON templates', 'level': 1, 'size': 16},
            {'text': '  - Boundary hierarchy definitions', 'level': 2, 'size': 15},
            {'text': '  - Localization messages', 'level': 2, 'size': 15},
            {'text': '  - Search queries', 'level': 2, 'size': 15},
            {'text': '', 'size': 14},
            {'text': 'Output: Generated artifacts', 'level': 1, 'size': 16},
            {'text': '  - IDs tracking file', 'level': 2, 'size': 15},
            {'text': '  - Excel templates', 'level': 2, 'size': 15},
            {'text': '  - Test reports', 'level': 2, 'size': 15},
            {'text': '', 'size': 14},
            {'text': 'Configuration Layer', 'size': 20, 'bold': True},
            {'text': '.env - Environment variables', 'level': 1, 'size': 16},
            {'text': 'pytest.ini - Test configuration', 'level': 1, 'size': 16},
        ],
        [
            {'text': 'Test Execution Flow', 'size': 20, 'bold': True},
            {'text': '1. Authentication', 'level': 0, 'size': 16},
            {'text': '2. Create hierarchy', 'level': 0, 'size': 16},
            {'text': '3. Validate creation', 'level': 0, 'size': 16},
            {'text': '4. Setup localization', 'level': 0, 'size': 16},
            {'text': '5. Generate template', 'level': 0, 'size': 16},
            {'text': '6. Poll generation status', 'level': 0, 'size': 16},
            {'text': '7. Download template', 'level': 0, 'size': 16},
            {'text': '8. Upload populated data', 'level': 0, 'size': 16},
            {'text': '9. Process data', 'level': 0, 'size': 16},
            {'text': '10. Validate processing', 'level': 0, 'size': 16},
            {'text': '11. Multi-lang validation', 'level': 0, 'size': 16},
            {'text': '12. Relationship search', 'level': 0, 'size': 16},
            {'text': '', 'size': 14},
            {'text': 'Reporting', 'size': 20, 'bold': True},
            {'text': 'HTML & Allure reports', 'level': 1, 'size': 16},
        ]
    )

    # Slide 6: Key Features - Part 1
    add_section_header(prs, "Key Features")

    add_content_slide(prs, "Core Testing Capabilities", [
        {'text': 'Sequential Test Execution', 'size': 20, 'bold': True},
        {'text': '15 interdependent tests with automatic dependency management', 'level': 1, 'size': 18},
        {'text': 'Tests must run in order (01 → 15) due to workflow dependencies', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Dynamic Hierarchy Generation', 'size': 20, 'bold': True},
        {'text': 'Each test run creates unique hierarchy types (e.g., TEST_D35387CC)', 'level': 1, 'size': 18},
        {'text': '7-level boundary structure: COUNTRY → PROVINCE → DISTRICT → POST ADMINISTRATIVE → LOCALITY → HEALTH FACILITY → VILLAGE', 'level': 1, 'size': 16},
        {'text': '', 'size': 18},
        {'text': 'ID Tracking & Cross-Test References', 'size': 20, 'bold': True},
        {'text': 'Generated IDs stored in output/ids.txt for reuse across tests', 'level': 1, 'size': 18},
        {'text': 'Ensures data consistency throughout test execution', 'level': 1, 'size': 18},
    ])

    # Slide 7: Key Features - Part 2
    add_content_slide(prs, "Advanced Features", [
        {'text': 'Template Automation', 'size': 20, 'bold': True},
        {'text': 'Automated template download from S3', 'level': 1, 'size': 18},
        {'text': 'Smart data population from reference samples', 'level': 1, 'size': 18},
        {'text': 'Header matching to ensure compatibility', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Multi-Language Localization Testing', 'size': 20, 'bold': True},
        {'text': 'Support for English, French, and Portuguese', 'level': 1, 'size': 18},
        {'text': 'Validation of localized boundary type names', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Asynchronous Operation Handling', 'size': 20, 'bold': True},
        {'text': 'Status polling for long-running operations', 'level': 1, 'size': 18},
        {'text': 'Automatic retry and wait mechanisms', 'level': 1, 'size': 18},
    ])

    # Slide 8: Key Features - Part 3
    add_content_slide(prs, "Quality & Reporting Features", [
        {'text': 'Comprehensive Reporting', 'size': 20, 'bold': True},
        {'text': 'HTML reports with detailed pass/fail summary', 'level': 1, 'size': 18},
        {'text': 'Allure reports with interactive dashboards', 'level': 1, 'size': 18},
        {'text': 'Execution time tracking and error details', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Reusable Components', 'size': 20, 'bold': True},
        {'text': 'APIClient wrapper with automatic authentication', 'level': 1, 'size': 18},
        {'text': 'Centralized configuration management', 'level': 1, 'size': 18},
        {'text': 'Modular payload loader', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Maintainability', 'size': 20, 'bold': True},
        {'text': 'Separation of test logic, payloads, and configuration', 'level': 1, 'size': 18},
        {'text': 'Easy to extend with new tests', 'level': 1, 'size': 18},
    ])

    # Slide 9: Demo
    add_section_header(prs, "Demo")

    add_content_slide(prs, "Test Execution Demo", [
        {'text': 'Running the complete test suite:', 'size': 20, 'bold': True},
        {'text': '', 'size': 18},
        {'text': 'Command: pytest tests/ -v', 'size': 18, 'bold': True},
        {'text': '', 'size': 18},
        {'text': 'What happens:', 'size': 20, 'bold': True},
        {'text': 'Test 01-02: Creates and validates 7-level boundary hierarchy', 'level': 1, 'size': 18},
        {'text': 'Test 03-04: Sets up and validates multi-language localization', 'level': 1, 'size': 18},
        {'text': 'Test 05-07: Generates and downloads boundary template', 'level': 1, 'size': 18},
        {'text': 'Test 08: Uploads populated boundary data', 'level': 1, 'size': 18},
        {'text': 'Test 09-11: Processes data and downloads results', 'level': 1, 'size': 18},
        {'text': 'Test 12-14: Validates localization in all languages', 'level': 1, 'size': 18},
        {'text': 'Test 15: Searches boundary relationships', 'level': 1, 'size': 18},
    ])

    # Slide 10: Demo Results
    add_content_slide(prs, "Demo Results", [
        {'text': 'Typical test run results:', 'size': 20, 'bold': True},
        {'text': '', 'size': 18},
        {'text': '14 tests passed', 'size': 22, 'bold': True},
        {'text': '1 test skipped (conditional on data availability)', 'size': 22, 'bold': True},
        {'text': '', 'size': 18},
        {'text': 'Generated artifacts:', 'size': 20, 'bold': True},
        {'text': 'output/ids.txt - Tracking file with all generated IDs', 'level': 1, 'size': 18},
        {'text': 'output/template_downloaded.xlsx - Downloaded template', 'level': 1, 'size': 18},
        {'text': 'output/sample_boundary.xlsx - Populated data for upload', 'level': 1, 'size': 18},
        {'text': 'reports/report.html - Comprehensive HTML report', 'level': 1, 'size': 18},
        {'text': 'allure-report/ - Interactive Allure dashboard', 'level': 1, 'size': 18},
    ])

    # Slide 11: Benefits
    add_section_header(prs, "Benefits")

    add_two_column_slide(
        prs,
        "Benefits & Value",
        [
            {'text': 'Time Savings', 'size': 20, 'bold': True},
            {'text': 'Automated 15-step workflow', 'level': 1, 'size': 16},
            {'text': 'Reduces hours of manual testing to minutes', 'level': 1, 'size': 16},
            {'text': 'One command execution: pytest tests/', 'level': 1, 'size': 16},
            {'text': '', 'size': 14},
            {'text': 'Quality Assurance', 'size': 20, 'bold': True},
            {'text': 'Consistent test execution', 'level': 1, 'size': 16},
            {'text': 'Comprehensive validation', 'level': 1, 'size': 16},
            {'text': 'Early bug detection', 'level': 1, 'size': 16},
            {'text': 'Reduced human error', 'level': 1, 'size': 16},
            {'text': '', 'size': 14},
            {'text': 'Maintainability', 'size': 20, 'bold': True},
            {'text': 'Modular architecture', 'level': 1, 'size': 16},
            {'text': 'Easy to extend and modify', 'level': 1, 'size': 16},
            {'text': 'Clear separation of concerns', 'level': 1, 'size': 16},
        ],
        [
            {'text': 'Reusability', 'size': 20, 'bold': True},
            {'text': 'Shared utility modules', 'level': 1, 'size': 16},
            {'text': 'Template-based payloads', 'level': 1, 'size': 16},
            {'text': 'Configurable via .env', 'level': 1, 'size': 16},
            {'text': '', 'size': 14},
            {'text': 'Visibility & Reporting', 'size': 20, 'bold': True},
            {'text': 'Multiple report formats', 'level': 1, 'size': 16},
            {'text': 'Detailed error tracking', 'level': 1, 'size': 16},
            {'text': 'Execution history', 'level': 1, 'size': 16},
            {'text': 'Test trends and analytics', 'level': 1, 'size': 16},
            {'text': '', 'size': 14},
            {'text': 'Scalability', 'size': 20, 'bold': True},
            {'text': 'Easy to add new tests', 'level': 1, 'size': 16},
            {'text': 'Supports CI/CD integration', 'level': 1, 'size': 16},
            {'text': 'Can handle complex workflows', 'level': 1, 'size': 16},
        ]
    )

    # Slide 12: Challenges
    add_section_header(prs, "Challenges & Learnings")

    add_content_slide(prs, "Challenges Faced", [
        {'text': 'Test Dependencies Management', 'size': 20, 'bold': True},
        {'text': 'Challenge: Tests depend on previous test outputs', 'level': 1, 'size': 18},
        {'text': 'Solution: Implemented ID tracking via output/ids.txt', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Case Sensitivity Issues', 'size': 20, 'bold': True},
        {'text': 'Challenge: API rejected mixed-case boundary types', 'level': 1, 'size': 18},
        {'text': 'Solution: Standardized to UPPERCASE naming convention', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Template Header Matching', 'size': 20, 'bold': True},
        {'text': 'Challenge: Uploaded data headers must match generated template', 'level': 1, 'size': 18},
        {'text': 'Solution: Created prepare_template_for_upload.py automation script', 'level': 1, 'size': 18},
    ])

    # Slide 13: More Challenges
    add_content_slide(prs, "Technical Challenges", [
        {'text': 'Asynchronous Operations', 'size': 20, 'bold': True},
        {'text': 'Challenge: File generation and processing take time', 'level': 1, 'size': 18},
        {'text': 'Solution: Implemented polling mechanism with status checks', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Excel File Operations', 'size': 20, 'bold': True},
        {'text': 'Challenge: Handling dynamic Excel templates', 'level': 1, 'size': 18},
        {'text': 'Solution: Used openpyxl for robust Excel manipulation', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Test Data Consistency', 'size': 20, 'bold': True},
        {'text': 'Challenge: Ensuring data flows correctly across all tests', 'level': 1, 'size': 18},
        {'text': 'Solution: Created reference sample file (never modified)', 'level': 1, 'size': 18},
    ])

    # Slide 14: Learnings
    add_content_slide(prs, "Key Learnings", [
        {'text': 'Best Practices Learned', 'size': 20, 'bold': True},
        {'text': '', 'size': 18},
        {'text': 'Modular design is crucial for maintainability', 'level': 0, 'size': 18},
        {'text': 'Separate configuration from code', 'level': 0, 'size': 18},
        {'text': 'Use environment variables for sensitive data', 'level': 0, 'size': 18},
        {'text': 'Implement comprehensive logging and reporting', 'level': 0, 'size': 18},
        {'text': 'Status polling is essential for async operations', 'level': 0, 'size': 18},
        {'text': 'Template automation saves significant time', 'level': 0, 'size': 18},
        {'text': 'Clear documentation prevents common mistakes', 'level': 0, 'size': 18},
        {'text': 'Sequential tests need careful dependency tracking', 'level': 0, 'size': 18},
    ])

    # Slide 15: Future Enhancements
    add_content_slide(prs, "Future Roadmap", [
        {'text': 'Planned Enhancements', 'size': 22, 'bold': True},
        {'text': '', 'size': 18},
        {'text': 'CI/CD Integration', 'size': 20, 'bold': True},
        {'text': 'Integrate with Jenkins/GitHub Actions', 'level': 1, 'size': 18},
        {'text': 'Automated test execution on code commits', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Parallel Test Execution', 'size': 20, 'bold': True},
        {'text': 'Run independent tests in parallel for faster execution', 'level': 1, 'size': 18},
        {'text': '', 'size': 18},
        {'text': 'Enhanced Reporting', 'size': 20, 'bold': True},
        {'text': 'Real-time test execution dashboard', 'level': 1, 'size': 18},
        {'text': 'Integration with test management tools', 'level': 1, 'size': 18},
    ])

    # Slide 16: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add "Thank You" text
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)

    thank_you_box = slide.shapes.add_textbox(left, top, width, height)
    thank_you_frame = thank_you_box.text_frame
    thank_you_frame.text = "Thank You!"

    p = thank_you_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add contact/repo info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    contact_frame = contact_box.text_frame
    contact_frame.text = "Repository: github.com/Shreya-egov/API_Automation"

    p = contact_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 102, 204)

    # Save presentation
    output_file = "API_Automation_Framework_Complete.pptx"
    prs.save(output_file)
    print(f"\nPresentation created successfully: {output_file}")
    print(f"\nSlides created:")
    print(f"  1. Title Slide")
    print(f"  2-3. Problem Statement")
    print(f"  4-5. Solution Overview")
    print(f"  6-7. Architecture Structure")
    print(f"  8-10. Key Features")
    print(f"  11-12. Demo")
    print(f"  13-14. Benefits")
    print(f"  15-17. Challenges & Learnings")
    print(f"  18. Future Roadmap")
    print(f"  19. Thank You")
    print(f"\nTotal slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
