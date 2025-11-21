#!/usr/bin/env python3
"""
Script to create a compact 5-slide PowerPoint presentation for the API Automation Framework
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for item in content_items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text = item
            level = 0

        p = text_frame.add_paragraph()
        p.text = text
        p.level = level
        p.space_before = Pt(6)

    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()

    # Slide 1: Title & Overview
    add_title_slide(
        prs,
        "API Automation Framework",
        "Boundary Management & Localization Services\n\n" +
        "15 Sequential Tests | End-to-End Workflow\n" +
        "Python • Pytest • Allure Reports • CI/CD Ready\n\n" +
        "Coverage: Hierarchy Creation, Localization, Template Generation,\n" +
        "File Operations, Data Processing & Validation"
    )

    # Slide 2: Section 1 - Intelligent API Automation with Python
    add_content_slide(
        prs,
        "1. Intelligent API Automation with Python",
        [
            "Technology Stack",
            ("Python 3.8+ with modern tooling", 1),
            ("Pytest - Testing framework with fixtures & plugins", 1),
            ("Requests - HTTP client for API calls", 1),
            ("Openpyxl - Excel file automation", 1),
            ("Allure - Interactive reporting", 1),
            "",
            "Smart Automation Features",
            ("Dynamic hierarchy generation (unique TEST_XXXXXXXX per run)", 1),
            ("Automated ID tracking across tests (output/ids.txt)", 1),
            ("Intelligent template handling with header matching", 1),
            ("Adaptive polling for async operations", 1),
            ("Multi-language localization (English, French, Portuguese)", 1),
            "",
            "Why Python?",
            ("Scalable, modular, readable, extensive libraries", 1)
        ]
    )

    # Slide 3: Section 2 - Framework Setup & Core Workflows
    add_content_slide(
        prs,
        "2. Framework Setup & Core Automation Workflows",
        [
            "Modular Architecture: Tests → Utils → API",
            ("api_client.py - HTTP wrapper | auth.py - OAuth2 tokens", 1),
            ("config.py - Environment vars | data_loader.py - JSON payloads", 1),
            "",
            "Quick Setup: Clone → Install → Configure .env → Run",
            ("pip install -r requirements.txt", 1),
            ("Configure BASE_URL, credentials, TENANTID in .env", 1),
            ("pytest tests/ -v --html=reports/report.html", 1),
            "",
            "Core Workflows (15 Tests)",
            ("Hierarchy: Create 7-level structure → Search & validate", 1),
            ("Localization: Upsert messages → Multi-language verification", 1),
            ("Templates: Generate → Download → Populate → Upload", 1),
            ("Processing: Process data → Validate relationships", 1),
            "",
            "Boundary Structure: COUNTRY → PROVINCE → DISTRICT → POST ADMINISTRATIVE → LOCALITY → HEALTH FACILITY → VILLAGE"
        ]
    )

    # Slide 4: Section 3 - Advanced Automation + AI-Driven Test Coverage
    add_content_slide(
        prs,
        "3. Advanced Automation + AI-Driven Test Coverage",
        [
            "Intelligent Template Automation",
            ("prepare_template_for_upload.py script", 1),
            ("Auto-download from S3 → Extract dynamic headers → Map columns", 1),
            ("Populate with reference data → Validate → Ready for upload", 1),
            ("Zero manual intervention required", 1),
            "",
            "Advanced Test Features",
            ("Cross-test dependencies via ID tracking", 1),
            ("Status polling with timeout (generate & process operations)", 1),
            ("Conditional test execution (skip when dependencies not met)", 1),
            ("Multi-language validation across all locales", 1),
            ("End-to-end relationship verification", 1),
            "",
            "Test Results & Coverage",
            ("14 tests passed, 1 skipped, execution: 80.10 seconds", 1),
            ("100% success rate | Complete workflow coverage", 1)
        ]
    )

    # Slide 5: Section 4 - Seamless CI/CD Integration
    add_content_slide(
        prs,
        "4. Seamless CI/CD Integration",
        [
            "CI/CD Ready Design",
            ("Environment-based config (.env) - no hardcoded credentials", 1),
            ("pytest.ini & requirements.txt for consistent setup", 1),
            ("Exit codes for pass/fail detection", 1),
            ("Artifact generation: HTML, Allure reports, logs", 1),
            "",
            "Platform Support",
            ("Jenkins: Jenkinsfile with stages (Setup → Test → Report)", 1),
            ("GitHub Actions: YAML workflow on push/PR", 1),
            ("GitLab CI: .gitlab-ci.yml with test & report stages", 1),
            "",
            "Continuous Testing Strategy",
            ("Triggers: Every commit, PR, scheduled nightly, on-demand", 1),
            ("Environments: Dev → QA → Staging → Production", 1),
            ("Notifications: Slack/Teams/Email on failures", 1),
            "",
            "Reports: HTML (self-contained) + Allure (interactive dashboard)",
            ("Open Allure: python3 -m http.server 8080 -d allure-report", 1)
        ]
    )

    # Save presentation
    prs.save('API_Automation_Framework_Presentation.pptx')
    print("✓ Compact presentation created successfully!")
    print("✓ File: API_Automation_Framework_Presentation.pptx")
    print("✓ Total slides: 5")
    print("\nStructure:")
    print("  1. Title & Overview")
    print("  2. Intelligent API Automation with Python")
    print("  3. Framework Setup & Core Workflows")
    print("  4. Advanced Automation + AI-Driven Testing")
    print("  5. Seamless CI/CD Integration")

if __name__ == "__main__":
    create_presentation()
