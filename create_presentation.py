#!/usr/bin/env python3
"""
Script to create a PowerPoint presentation for the API Automation Framework
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    return slide

def add_section_header(prs, title):
    """Add a section header slide"""
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for item in content_items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text = item
            level = 0

        p = text_frame.add_paragraph()
        p.text = text
        p.level = level
        p.space_before = Pt(6)

    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column content slide"""
    slide_layout = prs.slide_layouts[3]  # Two content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    left_shape = slide.placeholders[1]
    left_frame = left_shape.text_frame
    left_frame.clear()

    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.space_before = Pt(6)

    right_shape = slide.placeholders[2]
    right_frame = right_shape.text_frame
    right_frame.clear()

    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.space_before = Pt(6)

    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "API Automation Framework",
        "Boundary Management & Localization Testing\nPython • Pytest • Allure"
    )

    # Slide 2: Overview
    add_content_slide(
        prs,
        "Overview",
        [
            "Comprehensive Python-based API automation testing framework",
            "Focus: Boundary management and localization microservices",
            "Built with: Python 3.8+, Pytest, Allure Reports",
            "15 sequential tests covering complete workflow",
            "Automated template handling and data management",
            "Multiple reporting formats (HTML, Allure)"
        ]
    )

    # Slide 3: Key Features
    add_content_slide(
        prs,
        "Key Features",
        [
            "Sequential Test Execution",
            ("15 tests covering complete boundary management workflow", 1),
            "Template Automation",
            ("Automated template download, data population, and upload", 1),
            "Modularity",
            ("Reusable utilities for authentication, API calls, and data management", 1),
            "Maintainability",
            ("Separation of test logic, payloads, and configuration", 1),
            "Rich Reporting",
            ("HTML reports and interactive Allure reports", 1)
        ]
    )

    # Slide 4: Project Structure
    add_content_slide(
        prs,
        "Project Structure",
        [
            "tests/ - 15 sequential test modules",
            "utils/ - Reusable utility modules",
            ("api_client.py - HTTP client wrapper", 1),
            ("auth.py - Authentication token management", 1),
            ("config.py - Configuration loader", 1),
            ("data_loader.py - Payload loader", 1),
            "payloads/ - JSON payload templates",
            "output/ - Test outputs and generated files",
            "reports/ - HTML and Allure test reports"
        ]
    )

    # Slide 5: Architecture
    add_content_slide(
        prs,
        "Architecture & Design",
        [
            "Layered Architecture",
            ("Test Layer → Utility Layer → API Layer", 1),
            "Configuration Management",
            (".env file for environment-specific settings", 1),
            "Dynamic Data Generation",
            ("Unique hierarchy types per test run", 1),
            "ID Tracking System",
            ("output/ids.txt maintains cross-test references", 1),
            "Automated Template Workflow",
            ("Download → Populate → Upload → Process", 1)
        ]
    )

    # Slide 6: Test Suite Overview
    add_content_slide(
        prs,
        "Test Suite - 15 Sequential Tests",
        [
            "Tests 01-02: Boundary Hierarchy (Create & Search)",
            "Tests 03-04: Localization (Upsert & Search)",
            "Tests 05-06: Generate Template (Trigger & Status)",
            "Tests 07-08: File Operations (Download & Upload)",
            "Tests 09-10: Process Data (Trigger & Status)",
            "Test 11: Download Processed File",
            "Tests 12-14: Multi-language Localization",
            ("French, Portuguese, English", 1),
            "Test 15: Boundary Relationship Search"
        ]
    )

    # Slide 7: Boundary Hierarchy Structure
    add_content_slide(
        prs,
        "7-Level Boundary Hierarchy",
        [
            "Level 1: COUNTRY",
            ("Top-level boundary (e.g., Mozambique)", 1),
            "Level 2: PROVINCE",
            ("State/Province level", 1),
            "Level 3: DISTRICT",
            ("District/County level", 1),
            "Level 4: POST ADMINISTRATIVE",
            ("Administrative post", 1),
            "Level 5: LOCALITY",
            ("Local area", 1),
            "Level 6: HEALTH FACILITY",
            ("Healthcare facility location", 1),
            "Level 7: VILLAGE",
            ("Smallest boundary unit", 1)
        ]
    )

    # Slide 8: Test Execution Flow
    add_content_slide(
        prs,
        "Test Execution Flow",
        [
            "1. Create Hierarchy → Unique TEST_XXXXXXXX type generated",
            "2. Search Hierarchy → Verify creation successful",
            "3. Upsert Localization → Add multi-language support",
            "4. Generate Template → Trigger Excel template creation",
            "5. Poll Status → Wait for generation to complete",
            "6. Download Template → Fetch from S3 storage",
            "7. Populate Data → Use prepare_template_for_upload.py",
            "8. Upload File → Send populated template",
            "9. Process Data → Trigger boundary data processing",
            "10. Verify Results → Search and validate relationships"
        ]
    )

    # Slide 9: Setup & Configuration
    add_content_slide(
        prs,
        "Setup & Configuration",
        [
            "Prerequisites",
            ("Python 3.8+, pip, Git, openpyxl", 1),
            "Installation Steps",
            ("1. Clone repository", 1),
            ("2. Create virtual environment", 1),
            ("3. Install dependencies: pip install -r requirements.txt", 1),
            ("4. Configure .env file", 1),
            ".env Configuration",
            ("BASE_URL, USERNAME, PASSWORD, TENANTID", 1),
            ("CLIENT_AUTH_HEADER for OAuth", 1),
            ("SEARCH_LIMIT, SEARCH_OFFSET", 1)
        ]
    )

    # Slide 10: Template Automation
    add_content_slide(
        prs,
        "Template Automation",
        [
            "prepare_template_for_upload.py Script",
            "Workflow:",
            ("1. Download template from S3 (using FileStore ID)", 1),
            ("2. Load reference sample: utils/sample_boundary.xlsx", 1),
            ("3. Extract headers from downloaded template", 1),
            ("4. Copy data rows from reference sample", 1),
            ("5. Save as output/sample_boundary.xlsx", 1),
            ("6. Ready for Test 08 upload", 1),
            "Benefits:",
            ("✓ Automated header matching", 1),
            ("✓ Eliminates manual data entry", 1),
            ("✓ Ensures consistency", 1)
        ]
    )

    # Slide 11: Running Tests
    add_content_slide(
        prs,
        "Running Tests",
        [
            "Run All Tests:",
            ("pytest tests/ -v", 1),
            "Run with HTML Report:",
            ("pytest tests/ --html=reports/report.html", 1),
            "Run with Allure Report:",
            ("pytest tests/ --alluredir=allure-results", 1),
            ("allure generate allure-results --clean -o allure-report", 1),
            ("python3 -m http.server 8080 -d allure-report", 1),
            "Fresh Test Run:",
            ("rm -f output/ids.txt output/*.xlsx", 1),
            ("pytest tests/ -v", 1)
        ]
    )

    # Slide 12: Reporting Options
    add_content_slide(
        prs,
        "Reporting Options",
        [
            "HTML Reports",
            ("Self-contained HTML file", 1),
            ("Pass/Fail summary with execution time", 1),
            ("Detailed error traces", 1),
            "Allure Reports",
            ("Rich interactive web interface", 1),
            ("Test execution trends and history", 1),
            ("Detailed logs and attachments", 1),
            ("Graphs and charts for analysis", 1),
            "Output Files",
            ("output/ids.txt - Generated IDs", 1),
            ("Excel templates and processed files", 1)
        ]
    )

    # Slide 13: Test Results
    add_content_slide(
        prs,
        "Test Results & Metrics",
        [
            "Recent Test Run:",
            ("✓ 14 tests passed", 1),
            ("⊘ 1 test skipped (processed file check)", 1),
            ("⚠ 1 warning (HTTPS verification)", 1),
            ("⏱ Total execution: 80.10 seconds", 1),
            "Success Metrics:",
            ("100% pass rate for executed tests", 1),
            ("Complete boundary hierarchy creation", 1),
            ("Multi-language localization verified", 1),
            ("Template automation fully functional", 1),
            ("End-to-end workflow validated", 1)
        ]
    )

    # Slide 14: Best Practices
    add_content_slide(
        prs,
        "Best Practices",
        [
            "Always run tests in sequence (01 → 15)",
            "Clear output directory before fresh runs",
            ("Prevents stale data conflicts", 1),
            "Never modify utils/sample_boundary.xlsx",
            ("It's the reference template", 1),
            "Use prepare_template_for_upload.py for Test 08",
            ("Ensures proper header matching", 1),
            "Keep .env file secure and out of git",
            ("Contains sensitive credentials", 1),
            "Review output/ids.txt after test runs",
            ("Verify ID generation and tracking", 1)
        ]
    )

    # Slide 15: Troubleshooting
    add_content_slide(
        prs,
        "Common Issues & Solutions",
        [
            "Test 01 Fails: 'INVALID_HIERARCHY_DEFINITION'",
            ("→ Ensure UPPERCASE boundary types", 1),
            "Test 08 Skipped: 'Sample file not found'",
            ("→ Run prepare_template_for_upload.py first", 1),
            "Test 09 Fails: 'BOUNDARY_SHEET_HEADER_ERROR'",
            ("→ Re-download template and re-populate", 1),
            "Authentication Failure",
            ("→ Verify .env credentials and CLIENT_AUTH_HEADER", 1),
            "Module Import Errors",
            ("→ Activate virtual environment", 1),
            "Allure Report Java Error",
            ("→ Use Python HTTP server alternative", 1)
        ]
    )

    # Slide 16: Key Achievements
    add_content_slide(
        prs,
        "Key Achievements",
        [
            "✓ Complete end-to-end automation",
            ("15 sequential tests cover entire workflow", 1),
            "✓ Template automation eliminates manual work",
            ("Automated download, populate, upload process", 1),
            "✓ Multi-language support validated",
            ("English, French, Portuguese localization", 1),
            "✓ Robust error handling and reporting",
            ("Detailed logs and comprehensive reports", 1),
            "✓ Maintainable and scalable architecture",
            ("Modular design, easy to extend", 1)
        ]
    )

    # Slide 17: Future Enhancements
    add_content_slide(
        prs,
        "Future Enhancements",
        [
            "Parallel Test Execution",
            ("For independent test scenarios", 1),
            "CI/CD Integration",
            ("Jenkins, GitLab CI, GitHub Actions", 1),
            "Enhanced Logging",
            ("Request/Response logging for debugging", 1),
            "Test Data Management",
            ("Database integration for test data", 1),
            "Performance Testing",
            ("Load and stress testing capabilities", 1),
            "API Contract Testing",
            ("Schema validation with JSON Schema", 1)
        ]
    )

    # Slide 18: Thank You
    add_title_slide(
        prs,
        "Thank You!",
        "Questions?\n\nRepository: github.com/Shreya-egov/API_Automation"
    )

    # Save presentation
    prs.save('API_Automation_Framework_Presentation.pptx')
    print("✓ Presentation created successfully!")
    print("✓ File: API_Automation_Framework_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
