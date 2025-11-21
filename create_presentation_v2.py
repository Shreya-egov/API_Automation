#!/usr/bin/env python3
"""
Script to create an AI-focused PowerPoint presentation for the API Automation Framework
Following modern structure: Python Automation → Framework Setup → AI-Driven Testing → CI/CD
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    subtitle_shape.text = subtitle

    return slide

def add_section_header(prs, title):
    """Add a section header slide"""
    slide_layout = prs.slide_layouts[2]  # Section header layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for item in content_items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text = item
            level = 0

        p = text_frame.add_paragraph()
        p.text = text
        p.level = level
        p.space_before = Pt(6)

    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()

    # ====== TITLE & INTRO ======

    # Slide 1: Title Slide
    add_title_slide(
        prs,
        "Intelligent API Automation Framework",
        "Boundary Management & Localization Services\nPython • Pytest • AI-Driven Testing • CI/CD Ready"
    )

    # Slide 2: Framework Overview
    add_content_slide(
        prs,
        "API Automation for Boundary Management",
        [
            "Comprehensive automation framework for:",
            ("Boundary Hierarchy Creation (7-level structure)", 1),
            ("Localization Management (Multi-language support)", 1),
            ("Template Generation & Processing", 1),
            ("File Operations (Download/Upload)", 1),
            ("Data Validation & Relationship Verification", 1),
            "",
            "15 Sequential Tests | End-to-End Workflow Coverage",
            "Built with Modern Python | AI-Enhanced | CI/CD Ready"
        ]
    )

    # ====== SECTION 1: INTELLIGENT API AUTOMATION WITH PYTHON ======

    # Slide 3: Section Header
    add_section_header(prs, "1. Intelligent API Automation with Python")

    # Slide 4: Why Python for API Automation?
    add_content_slide(
        prs,
        "Why Python for API Automation?",
        [
            "Modern Python Tooling",
            ("Pytest - Powerful testing framework with fixtures", 1),
            ("Requests - Elegant HTTP library", 1),
            ("Python-dotenv - Environment management", 1),
            ("Openpyxl - Excel automation", 1),
            "",
            "Smart & Scalable Design",
            ("Modular architecture with reusable utilities", 1),
            ("Dynamic test data generation", 1),
            ("Automated template handling", 1),
            ("Rich reporting with Allure", 1)
        ]
    )

    # Slide 5: Intelligent Automation Features
    add_content_slide(
        prs,
        "Intelligent Automation Features",
        [
            "Dynamic Hierarchy Generation",
            ("Unique hierarchy types per run (TEST_XXXXXXXX)", 1),
            ("Automatic ID tracking across tests", 1),
            "",
            "Smart Template Automation",
            ("Auto-download from S3", 1),
            ("Intelligent header matching", 1),
            ("Automated data population from reference samples", 1),
            "",
            "Adaptive Polling & Validation",
            ("Status polling for async operations", 1),
            ("Multi-language localization verification", 1),
            ("End-to-end relationship validation", 1)
        ]
    )

    # Slide 6: Technology Stack
    add_content_slide(
        prs,
        "Technology Stack",
        [
            "Core Technologies",
            ("Python 3.8+ - Modern language features", 1),
            ("Pytest - Test execution & fixtures", 1),
            ("Requests - HTTP client", 1),
            ("Openpyxl - Excel operations", 1),
            "",
            "Reporting & Monitoring",
            ("Pytest-HTML - Self-contained HTML reports", 1),
            ("Allure - Interactive web reports", 1),
            ("Custom logging & ID tracking", 1),
            "",
            "Configuration Management",
            ("Python-dotenv - Environment variables", 1),
            ("JSON payload templates", 1)
        ]
    )

    # ====== SECTION 2: FRAMEWORK SETUP & CORE WORKFLOWS ======

    # Slide 7: Section Header
    add_section_header(prs, "2. Framework Setup & Core Automation Workflows")

    # Slide 8: Framework Architecture
    add_content_slide(
        prs,
        "Modular Framework Architecture",
        [
            "Layered Design",
            ("Tests → Utilities → API → Services", 1),
            "",
            "Core Modules",
            ("api_client.py - HTTP wrapper with auth", 1),
            ("auth.py - OAuth2 token management", 1),
            ("config.py - Environment configuration", 1),
            ("data_loader.py - JSON payload loading", 1),
            ("request_info.py - Request metadata builder", 1),
            "",
            "Separation of Concerns",
            ("Test logic separate from data", 1),
            ("Reusable utilities across tests", 1),
            ("Environment-specific configuration", 1)
        ]
    )

    # Slide 9: Project Structure
    add_content_slide(
        prs,
        "Clean Project Structure",
        [
            "tests/ - 15 sequential test modules",
            ("test_01_boundary_hierarchy_create.py", 1),
            ("test_02_boundary_hierarchy_search.py", 1),
            ("... through test_15", 1),
            "",
            "utils/ - Reusable utility modules",
            ("API client, auth, config, data loaders", 1),
            ("sample_boundary.xlsx - Reference template", 1),
            "",
            "payloads/ - JSON templates by category",
            ("boundary_hierarchy/, localization/, etc.", 1),
            "",
            "output/ - Generated files & ID tracking",
            "reports/ - HTML & Allure reports"
        ]
    )

    # Slide 10: Setup & Installation
    add_content_slide(
        prs,
        "Quick Setup & Installation",
        [
            "1. Clone Repository",
            ("git clone https://github.com/Shreya-egov/API_Automation.git", 1),
            ("cd API_Automation", 1),
            "",
            "2. Install Dependencies",
            ("pip install -r requirements.txt", 1),
            "",
            "3. Configure Environment",
            ("Create .env with BASE_URL, credentials, tenant info", 1),
            "",
            "4. Verify Setup",
            ("pytest tests/test_01_boundary_hierarchy_create.py -v", 1),
            "",
            "Ready to automate in < 5 minutes!"
        ]
    )

    # Slide 11: Core Automation Workflows
    add_content_slide(
        prs,
        "Core Automation Workflows",
        [
            "Workflow 1: Boundary Hierarchy Management",
            ("Create 7-level hierarchy → Search → Validate", 1),
            "",
            "Workflow 2: Multi-Language Localization",
            ("Upsert messages → Search by locale → Verify", 1),
            "",
            "Workflow 3: Template Generation & Processing",
            ("Generate → Poll status → Download → Populate → Upload", 1),
            "",
            "Workflow 4: Data Processing & Validation",
            ("Process uploaded data → Verify relationships → Download results", 1),
            "",
            "All workflows automated end-to-end!"
        ]
    )

    # Slide 12: 7-Level Boundary Hierarchy
    add_content_slide(
        prs,
        "7-Level Boundary Hierarchy",
        [
            "Hierarchical Structure",
            ("COUNTRY → PROVINCE → DISTRICT →", 1),
            ("POST ADMINISTRATIVE → LOCALITY →", 1),
            ("HEALTH FACILITY → VILLAGE", 1),
            "",
            "Example: Mozambique Structure",
            ("Mozambique → Tete → Cahora Bassa →", 1),
            ("Chitima → Chibagadigo → CS de Chirodze Ponte →", 1),
            ("Chissua sedew", 1),
            "",
            "Features",
            ("Dynamic generation with unique type codes", 1),
            ("Parent-child relationship validation", 1),
            ("Multi-language support at all levels", 1)
        ]
    )

    # Slide 13: Hands-On: Running Tests
    add_content_slide(
        prs,
        "Running Tests - Quick Guide",
        [
            "Run All Tests (Sequential)",
            ("pytest tests/ -v", 1),
            "",
            "Run with HTML Report",
            ("pytest tests/ --html=reports/report.html --self-contained-html", 1),
            "",
            "Run with Allure Report",
            ("pytest tests/ --alluredir=allure-results", 1),
            ("allure generate allure-results --clean -o allure-report", 1),
            ("python3 -m http.server 8080 -d allure-report", 1),
            "",
            "Fresh Test Run",
            ("rm -f output/ids.txt output/*.xlsx", 1),
            ("pytest tests/ -v", 1)
        ]
    )

    # ====== SECTION 3: ADVANCED AUTOMATION + AI-DRIVEN TEST COVERAGE ======

    # Slide 14: Section Header
    add_section_header(prs, "3. Advanced Automation + AI-Driven Test Coverage")

    # Slide 15: AI-Enhanced Testing Capabilities
    add_content_slide(
        prs,
        "AI-Enhanced Testing Capabilities",
        [
            "Dynamic Test Generation",
            ("Unique hierarchy types generated per run", 1),
            ("Randomized test data with consistent structure", 1),
            ("Automated ID generation and tracking", 1),
            "",
            "Smart Validations",
            ("Intelligent template header matching", 1),
            ("Automated data consistency checks", 1),
            ("Multi-language validation across locales", 1),
            ("Relationship verification using generated IDs", 1),
            "",
            "Adaptive Behavior",
            ("Polling with timeout for async operations", 1),
            ("Skip tests when dependencies not met", 1),
            ("Auto-retry mechanisms for transient failures", 1)
        ]
    )

    # Slide 16: Intelligent Template Automation
    add_content_slide(
        prs,
        "Intelligent Template Automation",
        [
            "prepare_template_for_upload.py",
            "",
            "Smart Workflow:",
            ("1. Fetch template from S3 using generated FileStore ID", 1),
            ("2. Extract dynamic headers (match current hierarchy)", 1),
            ("3. Load reference sample data", 1),
            ("4. Intelligently map columns", 1),
            ("5. Populate with consistent test data", 1),
            ("6. Validate before upload", 1),
            "",
            "Benefits:",
            ("✓ Zero manual intervention", 1),
            ("✓ Automatic header adaptation", 1),
            ("✓ Prevents data mismatch errors", 1)
        ]
    )

    # Slide 17: Advanced Test Features
    add_content_slide(
        prs,
        "Advanced Test Features",
        [
            "Cross-Test Dependencies",
            ("IDs tracked in output/ids.txt", 1),
            ("Seamless data flow between tests", 1),
            "",
            "Multi-Language Testing",
            ("English, French, Portuguese localization", 1),
            ("Automated locale switching", 1),
            ("Validation across all supported languages", 1),
            "",
            "File Operation Automation",
            ("S3 download/upload", 1),
            ("Excel file manipulation", 1),
            ("Template generation and processing", 1),
            "",
            "Comprehensive Validation",
            ("Schema validation, data integrity, relationships", 1)
        ]
    )

    # Slide 18: Test Coverage & Metrics
    add_content_slide(
        prs,
        "Test Coverage & Metrics",
        [
            "15 Sequential Tests - Complete Coverage",
            ("Hierarchy: Create, Search", 1),
            ("Localization: Upsert, Search (4 languages)", 1),
            ("Templates: Generate, Poll, Download", 1),
            ("Files: Upload, Process, Download", 1),
            ("Validation: Relationships, Data integrity", 1),
            "",
            "Recent Test Results",
            ("✓ 14 tests passed", 1),
            ("⊘ 1 test skipped (conditional)", 1),
            ("⏱ Execution time: 80 seconds", 1),
            ("✓ 100% success rate for executed tests", 1),
            ("✓ End-to-end workflow validated", 1)
        ]
    )

    # Slide 19: Reporting & Analytics
    add_content_slide(
        prs,
        "Rich Reporting & Analytics",
        [
            "HTML Reports",
            ("Self-contained, shareable reports", 1),
            ("Pass/Fail summary with timestamps", 1),
            ("Detailed error traces and logs", 1),
            "",
            "Allure Reports - Interactive Dashboard",
            ("Test execution trends over time", 1),
            ("Graphs and charts for analysis", 1),
            ("Detailed test steps and attachments", 1),
            ("Environment and system information", 1),
            "",
            "Custom Tracking",
            ("output/ids.txt - All generated IDs", 1),
            ("Excel files - Downloaded templates", 1),
            ("Logs - Detailed execution traces", 1)
        ]
    )

    # ====== SECTION 4: SEAMLESS CI/CD INTEGRATION ======

    # Slide 20: Section Header
    add_section_header(prs, "4. Seamless CI/CD Integration")

    # Slide 21: CI/CD Ready Architecture
    add_content_slide(
        prs,
        "CI/CD Ready Architecture",
        [
            "Framework Design for CI/CD",
            ("Environment-based configuration (.env)", 1),
            ("No hardcoded values or credentials", 1),
            ("Exit codes for pass/fail detection", 1),
            ("Artifact generation (reports, logs)", 1),
            "",
            "CI/CD Compatible Features",
            ("pytest.ini for consistent configuration", 1),
            ("requirements.txt for dependency management", 1),
            ("Clean test isolation", 1),
            ("Parallel execution support (future)", 1),
            "",
            "Report Integration",
            ("JUnit XML for Jenkins/GitLab", 1),
            ("Allure reports for visualization", 1)
        ]
    )

    # Slide 22: Jenkins Integration Example
    add_content_slide(
        prs,
        "Jenkins Pipeline Integration",
        [
            "pipeline {",
            ("agent any", 1),
            ("stages {", 1),
            ("  stage('Setup') {", 2),
            ("    steps {", 2),
            ("      sh 'pip install -r requirements.txt'", 3),
            ("    }", 2),
            ("  }", 2),
            ("  stage('Test') {", 2),
            ("    steps {", 2),
            ("      sh 'pytest tests/ --alluredir=allure-results'", 3),
            ("    }", 2),
            ("  }", 2),
            ("  stage('Report') {", 2),
            ("    steps {", 2),
            ("      allure includeProperties: false, results: [[path: 'allure-results']]", 3),
            ("    }", 2),
            ("  }", 2),
            ("}}", 1)
        ]
    )

    # Slide 23: GitHub Actions Integration
    add_content_slide(
        prs,
        "GitHub Actions Integration",
        [
            "name: API Automation Tests",
            ("on: [push, pull_request]", 1),
            ("jobs:", 1),
            ("  test:", 2),
            ("    runs-on: ubuntu-latest", 2),
            ("    steps:", 2),
            ("      - uses: actions/checkout@v3", 3),
            ("      - name: Set up Python", 3),
            ("        uses: actions/setup-python@v4", 3),
            ("        with: python-version: '3.10'", 4),
            ("      - name: Install dependencies", 3),
            ("        run: pip install -r requirements.txt", 3),
            ("      - name: Run tests", 3),
            ("        run: pytest tests/ --alluredir=allure-results", 3),
            ("      - name: Generate report", 3),
            ("        if: always()", 3),
            ("        run: allure generate allure-results", 3)
        ]
    )

    # Slide 24: GitLab CI Integration
    add_content_slide(
        prs,
        "GitLab CI/CD Integration",
        [
            "stages:",
            ("  - test", 1),
            ("  - report", 1),
            "",
            "test_api:",
            ("  stage: test", 1),
            ("  image: python:3.10", 1),
            ("  script:", 1),
            ("    - pip install -r requirements.txt", 2),
            ("    - pytest tests/ --alluredir=allure-results --junitxml=report.xml", 2),
            ("  artifacts:", 1),
            ("    reports:", 2),
            ("      junit: report.xml", 3),
            ("    paths:", 2),
            ("      - allure-results/", 3),
            ("      - reports/", 3)
        ]
    )

    # Slide 25: Continuous Testing Strategy
    add_content_slide(
        prs,
        "Continuous Testing Strategy",
        [
            "Trigger Points",
            ("On every commit to main branch", 1),
            ("On pull request creation/update", 1),
            ("Scheduled nightly runs", 1),
            ("On-demand manual triggers", 1),
            "",
            "Environment Strategy",
            ("Dev → QA → Staging → Production", 1),
            ("Environment-specific .env files", 1),
            ("Separate credentials per environment", 1),
            "",
            "Notification & Alerting",
            ("Slack/Teams integration for failures", 1),
            ("Email reports for stakeholders", 1),
            ("Dashboard for real-time monitoring", 1)
        ]
    )

    # Slide 26: Best Practices & Tips
    add_content_slide(
        prs,
        "Best Practices for Production",
        [
            "Test Execution",
            ("Always run tests in sequence (dependencies)", 1),
            ("Clear output directory before fresh runs", 1),
            ("Use virtual environments", 1),
            "",
            "Data Management",
            ("Never modify reference templates", 1),
            ("Keep .env files secure and out of git", 1),
            ("Review output/ids.txt after runs", 1),
            "",
            "CI/CD Best Practices",
            ("Use environment variables for secrets", 1),
            ("Archive test reports as artifacts", 1),
            ("Set appropriate timeouts", 1),
            ("Implement retry logic for flaky tests", 1)
        ]
    )

    # Slide 27: Troubleshooting Guide
    add_content_slide(
        prs,
        "Common Issues & Quick Fixes",
        [
            "Test Failures",
            ("Hierarchy creation fails → Check UPPERCASE naming", 1),
            ("Template upload fails → Run prepare script first", 1),
            ("Authentication fails → Verify .env credentials", 1),
            "",
            "CI/CD Issues",
            ("Build fails → Check Python version compatibility", 1),
            ("Tests timeout → Increase pipeline timeout settings", 1),
            ("Reports not generated → Verify allure installation", 1),
            "",
            "Debug Mode",
            ("pytest tests/ -vv -s --tb=long", 1),
            ("Enable detailed logging", 1),
            ("Review allure-results/ directory", 1)
        ]
    )

    # Slide 28: Key Achievements
    add_content_slide(
        prs,
        "Framework Achievements",
        [
            "✓ Complete End-to-End Automation",
            ("15 tests covering entire workflow", 1),
            "",
            "✓ Intelligent Template Processing",
            ("Zero manual intervention required", 1),
            "",
            "✓ Multi-Language Support",
            ("English, French, Portuguese validated", 1),
            "",
            "✓ CI/CD Ready",
            ("Jenkins, GitHub Actions, GitLab CI compatible", 1),
            "",
            "✓ Comprehensive Reporting",
            ("HTML + Allure for detailed insights", 1),
            "",
            "✓ 100% Test Success Rate",
            ("14/14 executed tests passing", 1)
        ]
    )

    # Slide 29: Future Roadmap
    add_content_slide(
        prs,
        "Future Enhancements",
        [
            "AI & Machine Learning",
            ("AI-powered test generation from API specs", 1),
            ("ML-based anomaly detection in responses", 1),
            ("Predictive test failure analysis", 1),
            "",
            "Performance & Scale",
            ("Parallel test execution", 1),
            ("Load and stress testing", 1),
            ("Performance regression detection", 1),
            "",
            "Advanced Features",
            ("API contract testing with JSON Schema", 1),
            ("Real-time monitoring dashboard", 1),
            ("Integration with APM tools", 1),
            ("Database validation automation", 1)
        ]
    )

    # Slide 30: Thank You
    add_title_slide(
        prs,
        "Thank You!",
        "Questions?\n\nRepository: github.com/Shreya-egov/API_Automation\n\nBuilt with Python • Pytest • AI-Driven • CI/CD Ready"
    )

    # Save presentation
    prs.save('API_Automation_Framework_Presentation.pptx')
    print("✓ Presentation created successfully!")
    print("✓ File: API_Automation_Framework_Presentation.pptx")
    print("✓ Total slides: 30")
    print("\nStructure:")
    print("  • Introduction (2 slides)")
    print("  • Section 1: Intelligent API Automation with Python (4 slides)")
    print("  • Section 2: Framework Setup & Core Workflows (7 slides)")
    print("  • Section 3: Advanced Automation + AI-Driven Testing (6 slides)")
    print("  • Section 4: Seamless CI/CD Integration (8 slides)")
    print("  • Conclusion (3 slides)")

if __name__ == "__main__":
    create_presentation()
