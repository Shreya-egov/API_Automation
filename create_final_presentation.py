from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
TITLE_COLOR = RGBColor(31, 78, 121)  # Dark Blue
ACCENT_COLOR = RGBColor(68, 114, 196)  # Medium Blue
TEXT_COLOR = RGBColor(64, 64, 64)  # Dark Gray
HIGHLIGHT_COLOR = RGBColor(237, 125, 49)  # Orange

def add_title_slide(prs):
    """Slide 1: Title Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add background shape
    left = Inches(0)
    top = Inches(2.5)
    width = Inches(10)
    height = Inches(2.5)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()

    # Add title
    left = Inches(1)
    top = Inches(2.8)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "API Automation Framework"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    left = Inches(1)
    top = Inches(4.2)
    width = Inches(8)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Boundary Management & Localization Testing"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """INTRODUCTION (1 minute):

Good morning/afternoon everyone. Today I'll be presenting our API Automation Framework that we've built for testing Boundary Management and Localization services.

WHAT TO SAY:
- "Thank you for joining me today"
- "Over the past few weeks/months, we've developed a comprehensive automation framework"
- "This framework has transformed how we test our boundary management system"
- "I'll walk you through what we had before, what we've automated, and the significant benefits we're seeing"
- "The presentation will take about 10 minutes, and I'm happy to take questions at the end"

TONE: Confident and enthusiastic. Set the stage for the value this automation brings."""

    return slide

def add_problems_slide(prs):
    """Slide 2: Problems Before Automation"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Challenges Before Automation"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Add content box
    left = Inches(0.8)
    top = Inches(1.3)
    width = Inches(8.4)
    height = Inches(5.5)

    # Problem 1
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "⏱️ Time-Consuming Manual Testing"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = HIGHLIGHT_COLOR
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "• 15 test scenarios executed manually = 2-3 hours per test cycle"
    p.font.size = Pt(16)
    p.level = 0
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = "• Testing across 3 languages (English, French, Portuguese) = triple effort"
    p.font.size = Pt(16)
    p.level = 0

    # Problem 2
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "❌ Human Error & Inconsistency"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = HIGHLIGHT_COLOR
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "• Manual data entry mistakes in Excel templates"
    p.font.size = Pt(16)
    p.level = 0
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = "• Missed test steps or incomplete validation"
    p.font.size = Pt(16)
    p.level = 0

    # Problem 3
    txBox = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "📊 Poor Test Coverage & Reporting"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = HIGHLIGHT_COLOR
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "• No systematic regression testing - only critical paths tested"
    p.font.size = Pt(16)
    p.level = 0
    p.space_after = Pt(4)

    p = tf.add_paragraph()
    p.text = "• Limited visibility into test results and trends"
    p.font.size = Pt(16)
    p.level = 0

    p = tf.add_paragraph()
    p.text = "• Difficult to track issues across complex multi-step workflows"
    p.font.size = Pt(16)
    p.level = 0

    # Problem 4
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.6))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "🔄 No Test Data Management"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = HIGHLIGHT_COLOR
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "• Test data conflicts when multiple testers work simultaneously"
    p.font.size = Pt(16)
    p.level = 0

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """PROBLEMS BEFORE AUTOMATION (2-3 minutes):

WHAT TO SAY:
"Before we built this automation framework, our testing process had several significant challenges."

PROBLEM 1 - TIME CONSUMING (30-45 seconds):
- "Our boundary management system has 15 complex test scenarios"
- "Each scenario involves multiple API calls - creating hierarchies, uploading files, processing data, and validating localization"
- "Manually executing all 15 tests took 2-3 hours per cycle"
- "Since we support 3 languages, we had to repeat tests for English, French, and Portuguese - essentially tripling our effort"
- "This meant regression testing before each release took an entire day"

PROBLEM 2 - HUMAN ERROR (30-45 seconds):
- "Manual testing is prone to mistakes"
- "We work with Excel files containing boundary data - manual data entry led to typos and formatting errors"
- "Testers sometimes skipped steps or forgot to validate certain fields"
- "This resulted in bugs slipping through to production"

PROBLEM 3 - POOR COVERAGE & REPORTING (45 seconds):
- "Due to time constraints, we only tested critical happy paths"
- "Edge cases and negative scenarios were often skipped"
- "We had no standardized reporting - test results were tracked in spreadsheets"
- "It was difficult to identify patterns or trends in failures"
- "The complex workflow - create hierarchy, generate template, upload file, process data - made it hard to pinpoint where issues occurred"

PROBLEM 4 - TEST DATA CONFLICTS (30 seconds):
- "Multiple testers working in the same environment caused data collisions"
- "One tester's data would interfere with another's tests"
- "We had no systematic way to generate unique test data"

TRANSITION:
"These challenges were costing us time, reducing quality, and creating frustration for the team. That's why we decided to invest in automation."
"""

    return slide

def add_automation_slide(prs):
    """Slide 3: What We Automated & Benefits"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Automated Functionalities & Benefits"
    p = tf.paragraphs[0]
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Left column - What we automated
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.8))
    tf = left_box.text_frame

    p = tf.add_paragraph()
    p.text = "✅ What We Automated"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(12)

    items = [
        "End-to-End Workflow (15 Tests)",
        "  • Boundary hierarchy creation (7-level structure)",
        "  • Multi-language localization (3 languages)",
        "  • File operations (generate, download, upload)",
        "  • Data processing & validation",
        "  • Relationship & search queries",
        "",
        "Framework Features",
        "  • Data-driven testing with JSON payloads",
        "  • Smart polling for async operations",
        "  • Cross-test data persistence (IDs tracking)",
        "  • OAuth2 authentication",
        "  • Dynamic unique test data generation",
    ]

    for item in items:
        p = tf.add_paragraph()
        if item.startswith("  •"):
            p.text = item[3:]
            p.font.size = Pt(13)
            p.level = 1
            p.space_after = Pt(3)
        elif item == "":
            p.text = ""
            p.font.size = Pt(6)
        else:
            p.text = item
            p.font.size = Pt(16)
            p.font.bold = True
            p.level = 0
            p.space_after = Pt(6)
            p.font.color.rgb = TEXT_COLOR

    # Right column - Benefits
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
    tf = right_box.text_frame

    p = tf.add_paragraph()
    p.text = "📈 Key Benefits"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(12)

    benefits = [
        ("⚡ 95% Time Savings", "2-3 hours → 5 minutes"),
        ("🎯 100% Consistency", "Zero human error, repeatable"),
        ("📊 Enhanced Reporting", "HTML + Allure + detailed logs"),
        ("🌍 Multi-language", "3 locales tested automatically"),
        ("🔄 Continuous Testing", "Run on every code change"),
        ("📈 Better Coverage", "All 15 scenarios + edge cases"),
        ("🔍 Early Bug Detection", "Issues caught before production"),
        ("👥 Team Productivity", "QA focuses on exploratory testing"),
    ]

    for title, desc in benefits:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.level = 0
        p.space_after = Pt(2)
        p.font.color.rgb = HIGHLIGHT_COLOR

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.level = 1
        p.space_after = Pt(8)

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """WHAT WE AUTOMATED & BENEFITS (3-4 minutes):

WHAT TO SAY - AUTOMATION COVERAGE (1.5-2 minutes):
"Let me walk you through what we've automated."

END-TO-END WORKFLOW:
- "We've automated 15 comprehensive test scenarios covering the entire boundary management lifecycle"
- "This includes creating a 7-level boundary hierarchy - from Country down to Village level"
- "We test multi-language localization across English, French, and Portuguese"
- "The framework handles complex file operations - generating Excel templates from the server, downloading them, uploading populated data, and validating processed results"
- "We test data processing, search queries, and boundary relationship validation"

FRAMEWORK CAPABILITIES:
- "The framework is data-driven - all test data is externalized in JSON files, making it easy to modify"
- "It includes smart polling for asynchronous operations - waiting intelligently for server-side file generation to complete"
- "We have cross-test data persistence - tests share IDs and data automatically through a tracking file"
- "OAuth2 authentication is handled automatically"
- "Each test run generates unique test data to prevent collisions, even when multiple testers run tests simultaneously"

WHAT TO SAY - BENEFITS (1.5-2 minutes):
"The benefits we're seeing are significant."

TIME SAVINGS:
- "The most obvious benefit - what took 2-3 hours manually now runs in 5 minutes"
- "That's a 95% time savings"
- "We can now run full regression tests multiple times per day instead of once per week"

QUALITY IMPROVEMENTS:
- "100% consistency - the tests run exactly the same way every time, eliminating human error"
- "We've dramatically improved test coverage - we now test all scenarios including edge cases"
- "Bugs are caught earlier in the development cycle, reducing production issues"

REPORTING & VISIBILITY:
- "We generate comprehensive reports in HTML and Allure formats"
- "Detailed logging captures every API request and response"
- "This visibility helps us quickly identify and diagnose issues"

TEAM IMPACT:
- "Our QA team can now focus on exploratory testing and complex scenarios"
- "We've integrated tests into our CI/CD pipeline for continuous testing"
- "The framework supports multi-language testing automatically, ensuring our international users have a quality experience"

TRANSITION:
"A key enabler in building this framework quickly was leveraging Claude AI..."
"""

    return slide

def add_claude_benefits_slide(prs):
    """Slide 4: Claude AI Benefits & Conclusion"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Claude AI: Accelerating Development"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    # Claude benefits section
    benefits_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.2))
    tf = benefits_box.text_frame

    p = tf.add_paragraph()
    p.text = "🤖 How Claude AI Helped"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(12)

    claude_benefits = [
        ("⚡ Rapid Framework Development", "Built complete framework in days vs. weeks - AI generated boilerplate code, utilities, and test structure"),
        ("🧠 Intelligent Code Generation", "Created data loaders, API clients, authentication modules with best practices built-in"),
        ("🔧 Complex Logic Handling", "AI helped implement smart polling, file operations, and cross-test data persistence"),
        ("📚 Learning Accelerator", "Explained pytest features, Allure integration, and Python best practices on-demand"),
        ("🐛 Debugging Partner", "Quickly identified issues, suggested fixes, and optimized code"),
    ]

    for title, desc in claude_benefits:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.level = 0
        p.space_after = Pt(4)
        p.font.color.rgb = HIGHLIGHT_COLOR

        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.level = 1
        p.space_after = Pt(10)

    # Key metrics box
    metrics_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(8.4), Inches(1.5))
    tf = metrics_box.text_frame

    p = tf.add_paragraph()
    p.text = "📊 Impact Summary"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "⏱️ 95% time reduction  |  🎯 15 automated scenarios  |  🌍 3 languages  |  📈 Zero manual errors  |  🚀 5-minute test runs"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = TEXT_COLOR

    # Conclusion
    conclusion_box = slide.shapes.add_textbox(Inches(1.5), Inches(6.5), Inches(7), Inches(0.8))
    tf = conclusion_box.text_frame

    p = tf.add_paragraph()
    p.text = "Result: Faster releases, higher quality, happier team ✅"
    p.font.size = Pt(22)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = ACCENT_COLOR

    # Speaker notes
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = """CLAUDE AI BENEFITS & CONCLUSION (2-3 minutes):

WHAT TO SAY - CLAUDE AI BENEFITS (1.5-2 minutes):
"I want to highlight how Claude AI was instrumental in building this framework quickly and efficiently."

RAPID DEVELOPMENT:
- "Typically, building a comprehensive automation framework like this would take weeks"
- "With Claude AI as a development partner, we built it in days"
- "Claude generated the initial boilerplate code, created utility modules, and set up the test structure"
- "This allowed me to focus on business logic and test scenarios rather than infrastructure"

INTELLIGENT CODE GENERATION:
- "Claude created sophisticated components like our API client wrapper, authentication manager, and data loaders"
- "These weren't just simple templates - they followed Python best practices and included error handling"
- "The AI understood our requirements and generated production-quality code"

COMPLEX LOGIC:
- "Some of our requirements were complex - like smart polling that waits for asynchronous file generation"
- "Claude helped implement the retry logic, timeout handling, and state management"
- "File operations - downloading templates, merging Excel data, uploading - were all complex tasks that Claude helped streamline"

LEARNING & DEBUGGING:
- "Claude acted as an on-demand tutor, explaining pytest features, Allure reporting integration, and Python libraries"
- "When I encountered issues, Claude helped debug by analyzing error messages and suggesting fixes"
- "It optimized code for better performance and readability"

PRODUCTIVITY MULTIPLIER:
- "Having an AI pair programmer dramatically accelerated development"
- "I could iterate faster, try different approaches, and build features I might not have attempted alone"
- "Claude handled the tedious parts so I could focus on the creative problem-solving"

WHAT TO SAY - IMPACT SUMMARY (30-45 seconds):
"Let me summarize the impact of this automation."
- "We've achieved 95% time reduction - from 2-3 hours down to 5 minutes"
- "15 comprehensive scenarios running automatically"
- "Testing across 3 languages with zero additional effort"
- "Eliminated manual testing errors completely"
- "We now have fast, reliable, repeatable tests"

CONCLUSION (30 seconds):
"The result is clear - we can release faster with higher confidence in quality"
"Our QA team is happier because they're doing more interesting work"
"Our developers catch issues earlier in the development cycle"
"And our end users - across English, French, and Portuguese speaking regions - get a more reliable product"

CLOSING:
"This framework represents the future of our testing strategy - automated, intelligent, and continuously improving"
"Thank you for your time. I'm happy to take any questions."

ANTICIPATED QUESTIONS & ANSWERS:
Q: How long did it take to build the framework?
A: "Approximately 3-4 days of active development with Claude AI's assistance, compared to 2-3 weeks if built traditionally"

Q: Can this be extended to other services?
A: "Absolutely. The framework is modular - we can add new test modules and payloads easily. We're planning to extend it to other microservices"

Q: What's the maintenance overhead?
A: "Very low. When APIs change, we just update the JSON payloads. The framework itself is stable and well-structured"

Q: How do you handle test data cleanup?
A: "Each test run generates unique IDs, so data doesn't collide. We can also add cleanup scripts if needed"

Q: Can non-technical people use this?
A: "They can view the reports and understand results. Writing new tests requires Python knowledge, but modifying test data in JSON files is accessible"
"""

    return slide

# Generate all slides
print("Creating presentation...")
add_title_slide(prs)
add_problems_slide(prs)
add_automation_slide(prs)
add_claude_benefits_slide(prs)

# Save presentation
output_file = "/home/shreya-kumar/API_Automation/API_Automation_Presentation_Final.pptx"
prs.save(output_file)
print(f"✅ Presentation created successfully: {output_file}")
print("\n📊 Presentation Structure:")
print("  Slide 1: Title - API Automation Framework")
print("  Slide 2: Challenges Before Automation")
print("  Slide 3: Automated Functionalities & Benefits")
print("  Slide 4: Claude AI Benefits & Conclusion")
print("\n💡 Speaker notes included for each slide (view in PowerPoint's Notes section)")
print("⏱️  Total presentation time: 8-10 minutes")
